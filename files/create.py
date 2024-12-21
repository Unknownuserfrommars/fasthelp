import os
from docx import Document
from pptx import Presentation


class create:
    @staticmethod
    def txt(filename: str, savepath: str, content: str = ""):
        """
        Creates a text file with specified content.

        Parameters:
        filename (str): The name of the file (e.g., 'example.txt').
        savepath (str): The directory to save the file in.
        content (str): Content to write into the file. Default is an empty string.

        Raises:
        ValueError: If the filename does not end with '.txt'.
        FileNotFoundError: If the savepath directory does not exist.
        """
        if not filename.endswith('.txt'):
            raise ValueError("The filename must end with '.txt'.")
        if not os.path.exists(savepath):
            raise FileNotFoundError(f"The directory '{savepath}' does not exist.")
        file_path = os.path.join(savepath, filename)
        try:
            with open(file_path, 'w') as f:
                f.write(content)
            print(f"File '{filename}' created successfully at '{savepath}'.")
        except Exception as e:
            raise RuntimeError(f"An error occurred while creating the file: {e}")

    @staticmethod
    def document(filename: str, savepath: str, title: str = None, headings: list = None, content: list = None):
        """
        Creates a Word document (.docx) with specified content.

        Parameters:
        filename (str): Name of the document file.
        savepath (str): Path to save the document.
        title (str): Title of the document.
        headings (list): List of headings.
        content (list): List of paragraphs to add as content.
        """
        doc = Document()
        if title:
            doc.add_heading(title, level=0)
        if headings:
            for i, heading in enumerate(headings, start=1):
                doc.add_heading(heading, level=i)
        if content:
            for paragraph in content:
                doc.add_paragraph(paragraph)
        doc.save(os.path.join(savepath, filename))

    @staticmethod
    def presentation(filename: str, savepath: str, total_slides: int = 1):
        """
        Creates a PowerPoint presentation with the specified number of slides.

        Parameters:
        filename (str): Name of the presentation file.
        savepath (str): Path to save the presentation.
        total_slides (int): Number of slides to create. Default is 1.
        """
        prs = Presentation()
        for _ in range(total_slides):
            prs.slides.add_slide(prs.slide_layouts[5])  # Add blank slides
        prs.save(os.path.join(savepath, filename))