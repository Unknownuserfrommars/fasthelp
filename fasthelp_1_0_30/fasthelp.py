import pandas as pd # type: ignore
import numpy as np
import pprint
import subprocess
import os
import warnings
import openai # type: ignore
from docx import Document
from pptx import Presentation
import string
import platform
import distro

class WillNotBeImplementedError(NotImplementedError):
    pass

class NoAPIEndpointError(Exception):
    """
    No API Endpoint
    """
    pass

class GetVersionInfo:
    def __init__(self):
        self._versionInfoMessage = "You are using FastHelp version 1.0.30, developed by JDT."
    def print(self):
        print(self._versionInfoMessage)
    def ret(self):
        return self._versionInfoMessage
    def pprint(self):
        pprint.pprint(self._versionInfoMessage)

class DataHelper:
    class new:
        def __init__(self, 
                    data_type, 
                    save_name: str):
            """
            You are using `DataHelper.new()`, developed by JDT, version 0.13.4.
            Creates a empty dataset (e.g.: `list`, `dict`, `pd.DataFrame`, etc.)

            Parameters:
            data_type (keyword): Enter a dataset type (e.g.: `list()`, `dict()`, `pd.DataFrame()`, etc.) (NOTE: When you enter the data_type, DO add brackets (for example, `list()` not `list`))
            save_name (`str`): The name of the dataset

            Returns:
            data_type: A dataset (type based on the input parameter 'data_type')
            """
            try:
                exec(f'{save_name} = {data_type}')
            except Exception as e:
                raise SyntaxError(f"Invalid Syntax, an error happened. Error message: {e}")
    class totype:
        def DataFrame(self, 
                      dataset) -> pd.DataFrame:
            """
            You are using `DataHelper.totype.DataFrame()`, developed by JDT, version 0.3.8.
            Change a dataset to a `pd.DataFrame` (Or `pd.Series`)
            
            Parameter:
            dataset (data set): The original dataset you want to change the type.

            Returns:
            `pd.DataFrame` or `pd.Series`
            """
            if not isinstance(dataset, (pd.DataFrame, pd.Series)):
                return pd.DataFrame(dataset)
            else:
                raise TypeError(f'The dataset you\'ve entered is a {type(dataset)} type. Which is currently not supported.')
            
        def list(self, 
                 dataset) -> list:
            """
            You are using `DataHelper.totype.list()`, developed by JDT, version 0.2.3.
            Change a dataset to a Python list.
            
            Parameter:
            dataset (data set): The original dataset you want to change the type.

            Returns:
            `list`: The list
            """
            if not isinstance(dataset, list):
                try:
                    return list(dataset)
                except Exception as e:
                    raise SyntaxError(f'An error occured while turning the dataset {dataset} to a list: {e}')
            else:
                raise TypeError(f'The dataset that you\'ve entered is a {type(dataset)} type. Which is not supported.')
            
        def dict(self, 
                 dataset) -> dict:
            """
            You are using `DataHelper.totype.dict()`, developed by JDT, version 0.2.4.
            Change a dataset to a Python dictionary.
            
            Parameter:
            dataset (data set): The original dataset you want to change the type.

            Returns:
            `dict`: The dictionary
            """
            if not isinstance(dataset, dict):
                try:
                    return dict(dataset)
                except Exception as e:
                    raise SyntaxError(f'An error occured while turning the dataset {dataset} to a dictionary: {e}')
            else:
                raise TypeError(f'The dataset that you\'ve entered is a {type(dataset)} type. Which is currently not supported.')
            
        def set(self, 
                dataset) -> set:
            """
            You are using `DataHelper.totype.set()`, developed by JDT, version 0.1.7.
            Change a dataset to a Python set
            """
            if not isinstance(dataset, set):
                try:
                    return set(dataset)
                except Exception as e:
                    raise SyntaxError(f'An error occured while turning the dataset `{dataset}` to a set: {e}')
            else:
                raise TypeError(f'The dataset that you\'ve entered is a {type(dataset)} type, which is currently not supported.')
        
        def npArray(self, 
                    dataset) -> np.ndarray:
            """
            You are using `DataHelper.totype.npArray()`, developed by JDT, version 0.5.0.
            Change a dataset to a numpy Array.
            
            Parameter:
            dataset (data set): The original dataset you want to change the type.

            Returns:
            `np.ndarray`: The numpy Array
            """
            if not isinstance(dataset, np.ndarray):
                try:
                    return np.array(dataset)
                except Exception as e:
                    raise SyntaxError(f'An error occured while turning the dataset {dataset} to a np.Array: {e}')
            else:
                raise TypeError(f'The dataset that you\'ve entered is a {type(dataset)} type. Which is currently not supported.')
        
    # Add more if needed
            
class future:
    class About:
        def __init__(self):
            self._about_future = "the `future` module contains many future implementations. Developed by JDT, version 0.4.1."
        def print(self):
            print(self._about_future)
        def ret(self):
            return self._about_future
        def pprint(self):
            pprint.pprint(self._about_future)

    class Full_Read_Csv:
        def __init__(self, **kwargs):
            self.FirstAppearVersion = "1.0.9"
            self.FirstImplementVersion = "1.3.0"  # Hopefully...
            return pd.read_csv(**kwargs)
        
    class Full_Read_Excel:
        def __init__(self, **kwargs):
            self.FirstAppearVersion = "1.0.10"
            self.FirstImplementVersion = "1.3.0"
            return pd.read_excel(**kwargs)
        
    class Full_Open_File:
        def __init__(self, file, **kwargs):
            self.FirstAppearVersion = '1.0.18'
            self.FirstImplementVersion = '1.3.0'
            with open(file, **kwargs) as f:
                return f
        
    class No_More_JDT:
        def __init__(self):
            """
            Congrats, you've found an easter egg!
            """
            self.FirstAppearVersion = "1.0.0"
            self.FirstImplementVersion = None
            self.Note = "JDT FOREVER!!!"
            return self.Note
    
    class Complete_LifeHelper:
        def __init__(self):
            self.FirstAppearVersion = "1.0.15"
            self.FirstImplementVersion = None
            self.Note = "JDT DOES NOT HELP YOUR LIFE!"
            return self.Note

class FileHelper:
    class Open:
        def __init__(self):
            pass
        
        def Opentxt(self, 
                    file: str, 
                    name=None):
            """
            You are using `FileHelper.Open.Opentxt()`, developed by JDT, version 0.3.3.
            Currently, it only supports opening a file with read-only mode. More modes will be added in the future.

            Parameters:
            file: The .txt file you want to open.
            name (Optional): return the file with a name.

            Returns:
            `TextIOWrapper`: An opened file
            """
            if name:
                exec(f"{name} = open({file})")
            else:
                return open(file)
        
        def txt(self, 
                file, 
                name=None):
            """
            Wrapper function for `FileHelper.Open.Opentxt()`, version 0.8.1.
            """
            if name:
                return FileHelper.Open.Opentxt(file=file, name=name)  # type: ignore # noqa: F821
            else:
                return FileHelper.Open.Opentxt(file=file)  # type: ignore # noqa: F821
            
        def Opencsv(self, 
                    filepath: str, 
                    sep: str=None, 
                    delimiter: str=None, 
                    header: str=None, 
                    index_col: int|str=None) -> pd.DataFrame:
            """
            You are using `FileHelper.Open.Opencsv()`, developed by JDT, version 0.9.19.
            Currently, it only supports few keyword arguments. More arguments will be added in the future.

            Parameters:
            file: The .csv file you want to open
            (Optional arguments:)
            sep: 
            delimiter:
            header:
            index_col:

            Returns:
            `pd.DataFrame`: Can be used with `DataHelper.totype`

            Example:
            ```py
            csv1 = DataHelper.totype.dict(FileHelper.Open.Opencsv("C:/Users/usr/Desktop/Folder/testcsv/csv1.csv"))
            ```
            """
            return pd.read_csv(filepath_or_buffer=filepath, sep=sep, delimiter=delimiter, header=header, index_col=index_col)
        
        def Openexcel(self, 
                      filename: str, 
                      sheet_name: list | int, 
                      index_col: int=None) -> pd.DataFrame:
            """
            You are using `FileHelper.Open.Openexcel()`, developed by JDT, version 0.7.3.

            Parameters:
            file: file path to the .xls/.xlsx file
            sheet_name: the name of the specific sheet(s)
            index_col (Optional): Column indexes

            Returns:
            pd.DataFrame
            """
            return pd.read_excel(filename, sheet_name, index_col=index_col)
        
        # Add more if needed
    class Create:
        def Createtxt(filename: str, savepath: str):
            raise NotImplementedError("This is simple to create but will not be created in this version.")
        
        def CreateDocument(filename: str, 
                           savepath: str, 
                           suffix: str=".docx", 
                           title: str=None, 
                           heading_1: str=None, 
                           heading_2: str=None, 
                           heading_3: str=None, 
                           heading_4: str=None, 
                           heading_5: str=None, 
                           heading_6: str=None, 
                           heading_7: str=None, 
                           heading_8: str=None, 
                           heading_9: str=None, 
                           table_rows: int=None,
                           table_cols: int=None, 
                           content: list=None):
            """
            You are using `FileHelper.Create.CreateDocument(), developed by JDT, version 0.6.3.

            Params:
            filename: Filename
            savepath: Path to save
            suffix: either '.docx' or '.doc'
            (Optional Params:)
            title: Main title
            heading_1 ~ heading_9: Adds a heading, level of heading depends on the int after the underscore in the param name.
            table_rows, table_cols: Creates a table, with `table_rows` number of rows and `table_cols` number of cols.
            content: The main content. MUST BE A LIST.
            """
            assert suffix in ['.docx', '.doc'], "The suffix must be either '.docx' or '.doc', default is '.docx'."
            doc = Document()
            doc.add_heading(title, 0)
            doc.add_heading(heading_1, level=1)
            doc.add_heading(heading_2, level=2)
            doc.add_heading(heading_3, level=3)
            doc.add_heading(heading_4, level=4)
            doc.add_heading(heading_5, level=5)
            doc.add_heading(heading_6, level=6)
            doc.add_heading(heading_7, level=7)
            doc.add_heading(heading_8, level=8)
            doc.add_heading(heading_9, level=9)
            doc.add_table(table_rows, table_cols)
            for paragraph in content:
                doc.add_paragraph(paragraph)
            doc.save(f'{savepath}/{filename}{suffix}')

        def CreatePresentation(filename: str, 
                               savepath: str, 
                               suffix: str='.pptx', 
                               total_slides: int=1):
            """
            You are using `FileHelper.Create.CreatePresentation()`, developed by JDT, version 0.3.7.
            """
            assert suffix in ['.pptx', '.ppt'], "The suffix must be either '.pptx' or '.ppt'. Default is '.pptx'."
            prs = Presentation()
            for slide in range(0, total_slides + 1):
                exec(f'slide_{slide + 1} = prs.slide_layouts[{slide}]')
            prs.save(f'{savepath}/{filename}{suffix}')
            
            

class StringHelper:
    def Upper(self, 
              string: str) -> str:
        """
        You are using `StringHelper.Upper()`, developed by JDT, version 0.4.2.

        Parameter:
        string (`str`): the original string

        Returns:
        `str`
        """
        return string.upper()
    
    def Lower(self, 
              string: str) -> str:
        """
        You are using `StringHelper.Lower()`, developed by JDT, version 0.4.3.

        Parameter:
        string (`str`): the original string

        Returns:
        `str`
        """
        return string.lower()
    
    def capitalize_letters(self, 
                           text: str, 
                           letters: list[str] | str = "") -> str:
        """
        You are using `StringHelper.capitalize_letters()`, developed by JDT and ElectroBL, version 0.2.21.

        Parameters:
        text (`str`): the original text
        letters (`list`: list of `str` or `str`): the letters you want to capitalize in the original string

        Returns:
        `str`
        """
        letters_set = set(letters)
        capitalized_text = ''.join([char.upper() if char in letters_set else char for char in text])
        return capitalized_text
    
    def lower_letters(self, 
                      text: str, 
                      letters: list[str] | str) -> str:
        """
        You are using `StringHelper.lower_letters()`, developed by JDT, version 0.1.10.

        Parameters:
        text (`str`): the original text
        letters (`list`: list of `str`): the letters you want to lower in the original string

        Returns:
        `str`
        """
        letters_set = set(letters)
        lowered_text = ''.join([char.lower() if char in letters_set else char for char in text])
        return lowered_text
    
    def lowercase_letters(self):
        """
        You are using `StringHelper.lowercase_letters()`, develoepd by JDT, version 0.2.4.
        """
        return string.ascii_lowercase
    
class CHelper:
    def __init__(self, file_path: str):
        """
        You are using `CHelper.__init__()`, developed by JDT, version 0.1.8.
        
        Parameter:
        file_path (`str`): The path to the C/C++ file.
        
        Returns:
        None

        Example:
        ```py
        from fasthelp import CHelper
        file_path = 'example.c'
        c_helper = CHelper(file_path)
        c_helper.open_file()
        content = c_helper.read_file()
        if content:
            print(content)
        c_helper.runfile()
        ```
        """
        self.file_path = file_path
        self.file_content = None
        self.binary_path = 'output_binary'
    
    def open_file(self):
        """
        You are using `CHelper.open_file()`, developed by JDT, version 0.2.4.
        
        Parameter:
        None
        
        Returns:
        None
        """
        try:
            with open(self.file_path, 'r') as file:
                self.file_content = file.read()
            print(f"File '{self.file_path}' opened successfully.")
        except FileNotFoundError:
            print(f"File '{self.file_path}' not found.")
        except IOError as e:
            print(f"Error opening file '{self.file_path}': {e}.")

    def read_file(self):
        """
        You are using `CHelper.read_file()`, developed by JDT, version 0.3.8.
        
        Parameter(s):
        None
        
        Return(s):
        file_content (`str` or `None`): The content of the file if available, otherwise `None`.
        """
        if self.file_content is not None:
            return self.file_content
        else:
            print("No file content available. Please open the file using Chelper.open_file() first.")
            return None

    def runfile(self):
        """
        You are using `CHelper.runfile()`, developed by JDT, version 0.8.4.
        
        Parameter:
        None
        
        Returns:
        None
        """
        file_extension = os.path.splitext(self.file_path)[1]
        compiler = None
        
        if file_extension == '.c':
            compiler = 'gcc'
        elif file_extension == '.cpp':
            compiler = 'g++'
        else:
            print("Unsupported file type. Please provide a C (.c) or C++ (.cpp) file.")
            return

        compile_command = [compiler, self.file_path, '-o', self.binary_path]
        try:
            subprocess.run(compile_command, check=True)
            print(f"Compilation successful. Running the binary '{self.binary_path}'...")
            run_command = [f'./{self.binary_path}']
            subprocess.run(run_command, check=True)
        except subprocess.CalledProcessError as e:
            print(f"An error occurred during compilation or execution: {e}")
        finally:
            # Clean up the binary file after execution
            if os.path.exists(self.binary_path):
                os.remove(self.binary_path)

class JavaHelper:
    def __init__(self, file_path: str):
        """
        You are using `JavaHelper.__init__()`, developed by JDT, version 0.4.2.
        
        Parameter:
        file_path (`str`): The path to the Java file.
        
        Returns:
        None
        
        Example:
        ```py
        file_path = 'Example.java'  # Replace with your Java file path
        java_helper = JavaHelper(file_path)
        java_helper.open_file()
        content = java_helper.read_file()
        if content:
            print(content)
        java_helper.run()
        ```
        """
        self.file_path = file_path
        self.file_content = None
        self.class_name = os.path.splitext(os.path.basename(file_path))[0]

    def open_file(self):
        """
        You are using `JavaHelper.open_file()`, developed by JDT, version 0.1.1.
        
        Parameter:
        None
        
        Returns:
        None
        """
        try:
            with open(self.file_path, 'r') as file:
                self.file_content = file.read()
            print(f"File '{self.file_path}' opened successfully.")
        except FileNotFoundError:
            print(f"File '{self.file_path}' not found.")
        except IOError:
            print(f"Error opening file '{self.file_path}'.")

    def read_file(self):
        """
        You are using `JavaHelper.read_file()`, developed by JDT, version 0.3.2.
        
        Parameter:
        None
        
        Returns:
        file_content (`str` or `None`): The content of the file if available, otherwise `None`.
        """
        if self.file_content is not None:
            return self.file_content
        else:
            print("No file content available. Please open the file first.")
            return None

    def run(self, clean_class=True):
        """
        You are using `JavaHelper.run()`, developed by JDT, version 0.5.2.
        
        Parameter:
        None
        
        Returns:
        None
        """
        compile_command = ['javac', self.file_path]
        if clean_class:
            try:
                subprocess.run(compile_command, check=True)
                print(f"Compilation successful. Running the class '{self.class_name}'...")
                run_command = ['java', self.class_name]
                subprocess.run(run_command, check=True)
            except subprocess.CalledProcessError as e:
                print(f"An error occurred during compilation or execution: {e}")
            finally:
                # Clean up the class file after execution
                class_file = self.class_name + '.class'
                if os.path.exists(class_file):
                    os.remove(class_file)
        else:
            try:
                subprocess.run(compile_command, check=True)
                print(f"Compilation successful. Running the class '{self.class_name}'...")
                run_command = ['java', self.class_name]
                subprocess.run(run_command, check=True)
            except subprocess.CalledProcessError as e:
                print(f"An error occurred during compilation or execution: {e}")

class JSHelper:
    def __init__(self, file_path):
        """
        You are using `JSHelper.__init__()`, developed by JDT, version 0.2.4.
        JSHelper requires Node.js installed.
        
        Parameter:
        file_path (`str`): The path to the JavaScript file.
        
        Returns:
        None
        
        Example:
        ```py
        file_path = 'example.js'  # Replace with your JavaScript file path
        js_helper = JSHelper(file_path)
        js_helper.open_file()
        content = js_helper.read_file()
        if content:
            print(content)
        js_helper.run()
        ```
        """
        self.file_path = file_path
        self.file_content = None

    def open_file(self):
        """
        You are using `JSHelper.open_file()`, developed by JDT, version 0.1.9.
        
        Parameter:
        None
        
        Returns:
        None
        """
        try:
            with open(self.file_path, 'r') as file:
                self.file_content = file.read()
            print(f"File '{self.file_path}' opened successfully.")
        except FileNotFoundError:
            print(f"File '{self.file_path}' not found.")
        except IOError:
            print(f"Error opening file '{self.file_path}'.")

    def read_file(self):
        """
        You are using `JSHelper.read_file()`, developed by JDT, version 0.3.1.
        
        Parameters:
        None
        
        Returns:
        file_content (`str` or `None`): The content of the file if available, otherwise `None`.
        """
        if self.file_content is not None:
            return self.file_content
        else:
            print("No file content available. Please open the file first.")
            return None

    def run(self):
        """
        You are using `JSHelper.run()`, developed by JDT, version 0.6.3.
        
        Parameter:
        None
        
        Returns:
        None
        """
        run_command = ['node', self.file_path]
        try:
            subprocess.run(run_command, check=True)
            print(f"Execution of '{self.file_path}' completed successfully.")
        except subprocess.CalledProcessError as e:
            print(f"An error occurred during execution: {e}")

class RHelper:
    def __init__(self, file_path):
        """
        You are using `RHelper.__init__()`, developed by JDT, version 0.8.0.
        
        Parameter:
        file_path (`str`): The path to the R script file.
        
        Returns:
        None
        
        Example:
        ```py
        file_path = 'example.R'  # Replace with your R script file path
        r_helper = RHelper(file_path)
        r_helper.open_file()
        content = r_helper.read_file()
        if content:
            print(content)
        r_helper.run()
        ```
        """
        self.file_path = file_path
        self.file_content = None

    def open_file(self):
        """
        You are using `RHelper.open_file()`, developed by JDT, version 0.2.9.
        
        Parameter:
        None
        
        Returns:
        None
        """
        try:
            with open(self.file_path, 'r') as file:
                self.file_content = file.read()
            print(f"File '{self.file_path}' opened successfully.")
        except FileNotFoundError:
            print(f"File '{self.file_path}' not found.")
        except IOError:
            print(f"Error opening file '{self.file_path}'.")

    def read_file(self):
        """
        You are using `RHelper.read_file()`, developed by JDT, version 0.3.5.
        
        Parameter:
        None
        
        Returns:
        file_content (`str` or `None`): The content of the file if available, otherwise `None`.
        """
        if self.file_content is not None:
            return self.file_content
        else:
            print("No file content available. Please open the file first.")
            return None

    def run(self):
        """
        You are using `RHelper.run()`, developed by JDT, version 0.3.1.
        
        Parameter:
        None
        
        Returns:
        None
        """
        run_command = ['Rscript', self.file_path]
        try:
            subprocess.run(run_command, check=True)
            print(f"Execution of '{self.file_path}' completed successfully.")
        except subprocess.CalledProcessError as e:
            print(f"An error occurred during execution: {e}")

# Add more if needed
# These are listed in the main class. Will be put into a subclass
# called CLangHelper (Computing Language Helper) in the future

class FolderHelper:
    def __init__(self, folder_path: str):
        """
        You are using `FolderHelper.__init__()`, developed by JDT, version 0.6.4.
        
        Parameter:
        folder_path (`str`): The path to the folder.
        
        Returns:
        None
        
        Example:
        ```py
        folder_path = 'C:\\Users\\usr\\Desktop\\some\\what\\folder'  # Replace with your folder path
        folder_helper = FolderHelper(folder_path)
        folder_helper.listdir()
        folder_helper.new_dir('new_folder')
        folder_helper.new_file(type='text', filename='new_file.txt')
        ```
        """
        self.folder_path = folder_path
        if not os.path.exists(folder_path):
            print(f"Folder '{folder_path}' does not exist.")
        else:
            print(f"Folder '{folder_path}' opened successfully.")

    def listdir(self):
        """
        You are using `FolderHelper.listdir()`, developed by JDT, version 0.4.3.
        
        Parameter:
        None
        
        Returns:
        List of files and directories in the folder (list of `str`)
        """
        try:
            items = os.listdir(self.folder_path)
            print(f"Contents of '{self.folder_path}':")
            for item in items:
                print(item)
            return items
        except FileNotFoundError:
            print(f"Folder '{self.folder_path}' not found.")
            return []
        except IOError as e:
            print(f"Error listing contents of '{self.folder_path}': {e}.")
            return []

    def new_dir(self, name: str):
        """
        You are using `FolderHelper.new_dir()`, developed by JDT, version 0.7.6.
        
        Parameter:
        name (`str`): The name of the new folder to create.
        
        Returns:
        None
        """
        new_folder_path = os.path.join(self.folder_path, name)
        try:
            os.makedirs(new_folder_path)
            print(f"Folder '{new_folder_path}' created successfully.")
        except OSError as e:
            print(f"Error creating folder '{new_folder_path}': {e}")

    def new_file(self, type: str='text', open_in_vscode: bool=False, filename: str='new_file'):
        """
        You are using `FolderHelper.new_file()`, developed by JDT and ElectroBL, version 0.16.3.
        
        Parameters:
        type (`str`): The type of file to create ('text'/'py'/'java'/'js'/'c'/'cpp'/'r'/'ruby'). Default is 'text'.
        open_in_vscode (`bool`): Whether to open the file in VSCode. Default is False.
        filename (`str`): The name of the file to create. Default is 'new_file'. (NOTE: DO NOT ADD SUFFIX! (`file`, not `file.txt`))
        
        Returns:
        None
        """
        extensions = {
            'text': '.txt',
            'py': '.py',
            'java': '.java',
            'js': '.js',
            'c': '.c',
            'cpp': '.cpp',
            'r': '.R',
            'ruby': '.rb'
        }
        
        if type not in extensions:
            print("Unsupported file type. Supported types are: 'text', 'py', 'java', 'js', 'c', 'cpp', 'r', 'ruby.")
            return

        file_extension = extensions[type]
        if not filename.endswith(file_extension):
            file_path = os.path.join(self.folder_path, filename + file_extension)
        else:
            file_path = os.path.join(self.folder_path, filename)
        
        try:
            with open(file_path, 'w') as file:
                file.write('')  # Create an empty file
            print(f"File '{file_path}' created successfully.")
            
            if open_in_vscode:
                subprocess.run(['code', file_path], check=True)
                print(f"File '{file_path}' opened in VSCode.")
        except IOError as e:
            print(f"Error creating file '{file_path}': {e}")
        except subprocess.CalledProcessError as e:
            print(f"Error opening file '{file_path}' in VSCode: {e}")

class IDEHelper:
    class PyCharm:
        def __init__(self, file=None):
            """
            You are using `IDEHelper.PyCharm.__init__()`, developed by JDT, version 0.4.21.
            
            Parameter:
            file (`str` or list of `str`, optional): The path to the Python file or a list of Python files to open in PyCharm.
            
            Returns:
            None
            
            Example:
            ```py
            pycharm = IDEHelper.PyCharm('path/to/your/file.py')
            pycharm = IDEHelper.PyCharm(['path/to/your/file1.py', 'path/to/your/file2.py'])
            pycharm = IDEHelper.PyCharm()
            ```
            """
            warnings.warn("This is not fully implemented. Lots of error may happen. Use at your own risk.", Warning)
            command = ['charm']
            if file:
                if isinstance(file, list):
                    command.extend([f for f in file if f.endswith(('.py', '.ipynb'))])
                elif file.endswith(('.py', '.ipynb')):
                    command.append(file)
            try:
                subprocess.run(command, check=True)
                print(f"PyCharm opened successfully with file(s): {file}")
            except subprocess.CalledProcessError as e:
                print(f"Error opening PyCharm: {e}")

    class IntelliJ:
        def __init__(self, file=None):
            """
            You are using `IDEHelper.IntelliJ.__init__()`, developed by JDT, version 0.4.9.
            
            Parameter:
            file (`str` or list of `str`, optional): The path to the Java/Kotlin file or a list of Java/Kotlin files to open in IntelliJ.
            
            Returns:
            None
            
            Example:
            ```py
            intellij = IDEHelper.IntelliJ('path/to/your/file.java')
            intellij = IDEHelper.IntelliJ(['path/to/your/file1.java', 'path/to/your/file2.kt'])
            intellij = IDEHelper.IntelliJ()
            ```
            """
            warnings.warn("This is not fully implemented. Lots of error may happen. Use at your own risk.", Warning)
            command = ['idea']
            if file:
                if isinstance(file, list):
                    command.extend([f for f in file if f.endswith(('.java', '.kt'))])
                elif file.endswith(('.java', '.kt')):
                    command.append(file)
            try:
                subprocess.run(command, check=True)
                print(f"IntelliJ opened successfully with file(s): {file}")
            except subprocess.CalledProcessError as e:
                print(f"Error opening IntelliJ: {e}")

    class VSCode:
        def __init__(self, file=None):
            """
            This is safe to use.
            Without specifying as 'safe-to-use', all other modules should be used at the users own risk.
            You are using `IDEHelper.VSCode.__init__()`, developed by JDT, version 0.4.5.
            
            Parameter:
            file (`str` or list of `str`, optional): The path to the file or a list of files to open in VSCode.
            
            Returns:
            None
            
            Example:
            ```py
            vscode = IDEHelper.VSCode('path/to/your/file.py')
            vscode = IDEHelper.VSCode(['path/to/your/file1.py', 'path/to/your/file2.py'])
            vscode = IDEHelper.VSCode()
            ```
            """
            command = ['code']
            if file:
                if isinstance(file, list):
                    command.extend(file)
                else:
                    command.append(file)
            try:
                subprocess.run(command, check=True)
                print(f"VSCode opened successfully with file(s): {file}")
            except subprocess.CalledProcessError as e:
                print(f"Error opening VSCode: {e}")

        def get_version(self):
            """
            You are using `IDEHelper.VSCode.get_version()`, developed by JDT, version 0.6.6.
            
            Parameter:
            None
            
            Returns:
            `str`: Version of VSCode
            """
            try:
                result = subprocess.run(['code', '--version'], check=True, stdout=subprocess.PIPE, text=True)
                version_info = result.stdout.splitlines()[0]
                return f"VSCode {version_info}"
            except subprocess.CalledProcessError as e:
                return f"Error retrieving VSCode version: {e}"

    class VisualStudio_2019:
        def __init__(self, file=None):
            """
            You are using `IDEHelper.VisualStudio_2019.__init__()`, developed by JDT, version 0.4.4.
            
            Parameter:
            file (`str` or list of `str`, optional): The path to the C/C++ file or a list of C/C++ files to open in Visual Studio 2019.
            
            Returns:
            None
            
            Example:
            ```py
            vs2019 = IDEHelper.VisualStudio_2019('path/to/your/file.cpp')
            vs2019 = IDEHelper.VisualStudio_2019(['path/to/your/file1.c', 'path/to/your/file2.cpp'])
            vs2019 = IDEHelper.VisualStudio_2019()
            ```
            """
            warnings.warn("This is not fully implemented. Lots of error may happen. Use at your own risk.", Warning)
            command = ['devenv']
            if file:
                if isinstance(file, list):
                    command.extend([f for f in file if f.endswith(('.c', '.cpp'))])
                elif file.endswith(('.c', '.cpp')):
                    command.append(file)
            try:
                subprocess.run(command, check=True)
                print(f"Visual Studio 2019 opened successfully with file(s): {file}")
            except subprocess.CalledProcessError as e:
                print(f"Error opening Visual Studio 2019: {e}")

    class RStudio:
        def __init__(self, file=None):
            """
            You are using `IDEHelper.RStudio.__init__()`, developed by JDT, version 0.3.14.
            
            Parameter:
            file (`str` or list of `str`, optional): The path to the R/RMarkdown file or a list of R/RMarkdown files to open in RStudio.
            
            Returns:
            None
            
            Example:
            ```py
            rstudio = IDEHelper.RStudio('path/to/your/file.R')
            rstudio = IDEHelper.RStudio(['path/to/your/file1.R', 'path/to/your/file2.Rmd'])
            rstudio = IDEHelper.RStudio()
            ```
            """
            warnings.warn("This is not fully implemented. Lots of error may happen. Use at your own risk.", Warning)
            command = ['rstudio']
            if file:
                if isinstance(file, list):
                    command.extend([f for f in file if f.endswith(('.R', '.Rmd'))])
                elif file.endswith(('.R', '.Rmd')):
                    command.append(file)
            try:
                subprocess.run(command, check=True)
                print(f"RStudio opened successfully with file(s): {file}")
            except subprocess.CalledProcessError as e:
                print(f"Error opening RStudio: {e}")

    class PyIDLE:
        def __init__(self, file=None):
            """
            This is safe to use to an extent. (NOTE: TO AN EXTENT!!!)
            Without specifying as 'safe-to-use', all other modules should be used at the users own risk.
            You are using `IDEHelper.PyIDLE.__init__()`, developed by JDT, version 0.10.5.
            
            Parameter(s):
            file (`str` or list of `str`, optional): The path to the Python file or a list of Python files to open in IDLE.
            
            Returns:
            None
            
            Example:
            ```py
            idle = IDEHelper.PyIDLE('path/to/your/file.py')
            idle = IDEHelper.PyIDLE(['path/to/your/file1.py', 'path/to/your/file2.py'])
            idle = IDEHelper.PyIDLE()
            ```
            """
            warnings.warn("This is not FULLY implemented. Lots of error may happen. Use at your own risk.", Warning)
            command = ['python']
            if file:
                if isinstance(file, list):
                    command.extend([f for f in file if f.endswith('.py')])
                elif file.endswith('.py'):
                    command.append(file)
            try:
                subprocess.run(command, check=True)
                print(f"IDLE opened successfully with file(s): {file}")
            except subprocess.CalledProcessError as e:
                print(f"Error opening IDLE: {e}")

        def get_version(self):
            """
            You are using `IDEHelper.PyIDLE.get_version()`, developed by JDT, version 0.2.4.
            
            Parameter:
            None
            
            Returns:
            `str`: Version of IDLE
            """
            try:
                result = subprocess.run(['python', '--version'], check=True, stdout=subprocess.PIPE, text=True)
                version_info = result.stdout.strip()
                return f"Python {version_info} Shell"
            except subprocess.CalledProcessError as e:
                return f"Error retrieving IDLE version: {e}"
            
class LifeHelper:
    def __init__(self):
        """
        Developed by JDT. Version 0.1.0.

        Well no. We aren't going to help u with ur life!
        """
        raise WillNotBeImplementedError("Will not be implemented. NEVER!!!")
    
    # Nah. Just don't think about it. Pls, this counts as an easter egg!
    # Easter egg #0002. First Released in: 1.0.15.

class GPTHelper:
    def __init__(self, api_key):
        self.api_key = api_key
        openai.api_key = self.api_key
        self.model = None
        self.latest_response = None

    def login(self, username, password):
        # Since OpenAI API uses API keys, we assume login is managed externally
        # and focus on setting the API key.
        # Can integrate an external auth system if needed.
        print("Login.")
        raise NotImplementedError("Not Implemented...")

    def select_model(self, model):
        self.model = model
        return f"Model {model} selected."

    def ask(self, prompt):
        if not self.model:
            raise ValueError("No model selected. Please select a model first.")
        response = openai.chat.completions.create(
            model=self.model,
            messages=[
                {"role": "system", "content": "You are a helpful assistant."},
                {"role": "user", "content": prompt},
            ]
        )
        self.latest_response = response
        return response

    def get_answer(self):
        if not self.latest_response:
            raise ValueError("No prompt asked. Please use the 'ask' function first.")
        return self.latest_response['choices'][0]['message']['content'] # type: ignore

    def subscribe(self):
        # OpenAI API does not have a subscribe endpoint, this is a placeholder
        raise NoAPIEndpointError("OpenAI API does not have a subscribe endpoint")

    def use_gpts(self, GPTs_name=None):
        # This is a placeholder for a custom endpoint to use specific GPTs
        # return f"Using GPTs: {GPTs_name} :)"
        raise NotImplementedError("Not Implemented...")

    def check(self, account=None):
        # This is a placeholder for checking account details
        # return f"Account {account} is active :)"
        raise NotImplementedError("Not Implemented...")

    def delete_account(self):
        # This is a placeholder as OpenAI does not have a delete account endpoint
        raise NoAPIEndpointError("OpenAI does not have a delete account endpoint.")

# I know this is NOT PEP 8 but the imports are just too much!
import psutil # noqa: E402
import socket # noqa: E402
import shutil # noqa: E402
import time # noqa: E402

class OSHelper:
    @staticmethod
    def get_os_version(linux_pretty: bool=False):
        '''
        You are using `OSHelper.get_os_version()`, developed by JDT, version 0.3.16.
        Params:
        linux_pretty (bool): defaults False, if True, sets the `pretty` param of distro.name() and distro.version() as True
        '''
        os_name = platform.system()
        if os_name == "Windows":
            ver = platform.release()
            ver_detail = platform.version()
            os_ver = f'{os_name}{ver}({ver_detail})'
        elif os_name == "Darwin":
            mac_ver = platform.mac_ver()[0]
            os_ver = f'macOS {mac_ver}'
        elif os_ver == 'Linux':
            dist_name = distro.name(linux_pretty)
            ver = distro.version(linux_pretty)
            os_ver = f'{dist_name}, {ver}'
        else:
            os_ver = """Unknown OS. Or this OS is not supported.
                     Current support OS are Windows, MacOS and Linux
                     If you want a new OS support, please report it at the Github repo
                     https://github.com/Unknownuserfrommars/fasthelp/issues
                     """
        return os_ver

    @staticmethod
    def get_system_architechure():
        return platform.architecture()[0]

    @staticmethod
    def get_cpu_info():
        return {
            "processor": platform.processor(),
            "cores": os.cpu_count(),
        }
    
    @staticmethod
    def get_memory_usage():
        virtual_mem = psutil.virtual_memory()
        return {
            'total': virtual_mem.total,
            'available': virtual_mem.available,
            'used': virtual_mem.used,
            'free': virtual_mem.free
        }

    @staticmethod
    def get_disk_space(path: str="/"):
        total, used, free = shutil.disk_usage(path)
        return {
            'total': total,
            'used': used,
            'free': free,
        }

    @staticmethod
    def get_network_info():
        hostname = socket.gethostname()
        ip_address = socket.gethostbyname(hostname)
        return {
            "hostname": hostname,
            "ip_address": ip_address
        }

    @staticmethod
    def get_system_uptime():
        uptime_s = time.time() - psutil.boot_time()
        uptime_str = time.strftime("%H:%M:%S", time.gmtime(uptime_s))
        return uptime_str

    @staticmethod
    def get_running_processes():
        DataHelper.new(list, 'processes')
        for proc in psutil.process_iter(['pid', 'name']):
            processes.append(proc.info)  # type: ignore # noqa: F821
        return processes  # type: ignore # noqa: F821\

    @staticmethod
    def get_computer_info():
        cpu_info = OSHelper.get_cpu_info()
        gpu_info = [{'info': 'GPU info not supported in this version. The reason is because JDT haven\'t found a perfect package for all the GPUs. NVIDIA is easy to work, but the others? Maybe not.'}]
        mem_info = OSHelper.get_memory_usage()
        return {
            "cpu": cpu_info,
            'gpu': gpu_info,
            'ram': mem_info,
        }
    
    @staticmethod
    def bundle(computer_info: bool=False):
        x = [OSHelper.get_computer_info(), OSHelper.get_system_architechure(), OSHelper.get_cpu_info(), OSHelper.get_memory_usage(), OSHelper.get_disk_space(),
                OSHelper.get_network_info(), OSHelper.get_system_uptime(), OSHelper.get_running_processes()]
        x = x.append(OSHelper.get_computer_info()) if computer_info else None
        return x